# ================================
# IntelliBot FINAL Backend
# OS + COA + DBMS + CD + Teacher Docs
# Gemini + FAISS + Multi Subject
# ================================

from fastapi import FastAPI, UploadFile, File
from pydantic import BaseModel
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
import os, json, uuid
from typing import List, Dict, Any

# Document readers
import fitz
from docx import Document as DocxDocument
from pptx import Presentation

# ================================
# FastAPI Setup
# ================================
app = FastAPI()

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)
app.mount("/uploads", StaticFiles(directory=UPLOAD_DIR), name="uploads")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)

# ================================
# LangChain Imports
# ================================
from langchain.schema import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain_google_genai import ChatGoogleGenerativeAI

# ================================
# Configuration
# ================================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

JSON_FILES = {
    "os": os.path.join(BASE_DIR, "os.json"),
    "coa": os.path.join(BASE_DIR, "coa.json"),
    "dbms": os.path.join(BASE_DIR, "dbms.json"),
    "cd": os.path.join(BASE_DIR, "cd.json"),
    "python": os.path.join(BASE_DIR, "python.json")   # NEW
}


EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"

# Gemini API Key
os.environ["GOOGLE_API_KEY"] = "AIzaSyBftvcbwNQbBszpt7iY-wYnDF_B512KuR8"

# ================================
# Models
# ================================
class QuestionRequest(BaseModel):
    topic: str
    question: str

# ================================
# Globals
# ================================
vectorstores = {}
qas = {}
chat_memory = []
MAX_MEMORY = 6
notes_db = {}

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
NOTES_FILE = os.path.join(BASE_DIR, "notes_data.json")

# Load notes from file at startup
# Load notes from disk when server starts
if os.path.exists(NOTES_FILE):
    with open(NOTES_FILE, "r") as f:
        notes_db = json.load(f)
else:
    notes_db = {}



embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

# ================================
# Initialize Gemini
# ================================
def init_llm():
    print("✅ Using Gemini")
    return ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.2)

llm = init_llm()

# ================================
# JSON → Documents
# ================================
def json_to_documents(json_data: Dict[str, Any]) -> List[Document]:
    docs = []

    for module in json_data.get("modules", []):
        for topic in module.get("topics", []):
            topic_name = topic.get("topic_name", "topic")
            content = topic.get("content", {})

            parts = []

            for key, value in content.items():

                # 🔥 HANDLE LISTS SAFELY
                if isinstance(value, list):
                    safe_list = []
                    for item in value:
                        if isinstance(item, dict):
                            safe_list.append(" ".join(str(v) for v in item.values()))
                        else:
                            safe_list.append(str(item))
                    value = " | ".join(safe_list)

                # 🔥 HANDLE DICTS
                elif isinstance(value, dict):
                    value = " ".join(str(v) for v in value.values())

                else:
                    value = str(value)

                parts.append(f"{key.title()}: {value}")

            text = "\n".join(parts)

            docs.append(
                Document(
                    page_content=text,
                    metadata={
                        "source": "syllabus",
                        "topic": topic_name
                    }
                )
            )

    return docs


# ================================
# Document readers
# ================================
def extract_text_from_pdf(path):
    text = ""
    doc = fitz.open(path)
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_docx(path):
    doc = DocxDocument(path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

# ================================
# Build vectorstore per subject
# ================================
def build_vectorstore_for_subject(subject, json_path):
    global vectorstores, qas

    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    docs = json_to_documents(data)

    splitter = RecursiveCharacterTextSplitter(chunk_size=700, chunk_overlap=120)
    chunks = splitter.split_documents(docs)

    vectorstore = FAISS.from_documents(chunks, embeddings)
    retriever = vectorstore.as_retriever(search_kwargs={"k": 3})

    qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=retriever)

    vectorstores[subject] = vectorstore
    qas[subject] = qa_chain

    print(f"✅ {subject} ready with {len(chunks)} chunks")

# ================================
# Add teacher docs into FAISS
# ================================
def add_teacher_document(path, subject, filename):

    if subject not in vectorstores:
        print("⚠ Subject not found for embedding:", subject)
        return

    if filename.endswith(".pdf"):
        text = extract_text_from_pdf(path)
    elif filename.endswith(".docx"):
        text = extract_text_from_docx(path)
    elif filename.endswith(".pptx"):
        text = extract_text_from_pptx(path)
    elif filename.endswith(".txt"):
        text = extract_text_from_txt(path)
    else:
        print("Unsupported file type")
        return

    doc = Document(page_content=text, metadata={"source":"teacher_note","subject":subject,"file":filename})

    splitter = RecursiveCharacterTextSplitter(chunk_size=700, chunk_overlap=150)
    chunks = splitter.split_documents([doc])

    vectorstores[subject].add_documents(chunks)

    print(f"📘 Added teacher document: {filename}")

# ================================
# ASK API
# ================================
@app.post("/ask")
async def ask_question(req: QuestionRequest):

    subject = req.topic.lower().strip()

    subject_map = {

    "operating systems": "os",
    "os": "os",
    "computer organization": "coa",
    "coa": "coa",
    "dbms": "dbms",
    "database": "dbms",
    "compiler design": "cd",
    "cd": "cd",
    "python": "python",
    "py": "python"
}


    subject_key = subject_map.get(subject, subject)

    if subject_key not in qas:
        return {"answer":"Subject not available"}

    qa = qas[subject_key]

    chat_memory.append(f"User: {req.question}")
    memory_context = "\n".join(chat_memory[-MAX_MEMORY:])

    prompt = f"""
You are a Computer Science academic assistant.

Conversation:
{memory_context}

Answer clearly from syllabus + teacher materials:
{req.question}
"""

    result = qa.run(prompt)

    chat_memory.append(f"Bot: {result}")
    return {"answer": result}

# ================================
# Notes upload + embedding
# ================================
@app.post("/notes/upload/{subject}")
async def upload_note(subject: str, file: UploadFile = File(...)):

    subject = subject.lower().strip()

    subject_map = {
        "operating systems": "os",
        "os": "os",
        "computer organization": "coa",
        "coa": "coa",
        "dbms": "dbms",
        "database": "dbms",
        "compiler design": "cd",
        "cd": "cd",
        "python": "python",
        "py": "python"
    }

    subject_key = subject_map.get(subject, subject)

    filename = f"{uuid.uuid4()}_{file.filename}"
    path = os.path.join(UPLOAD_DIR, filename)

    with open(path, "wb") as f:
        f.write(await file.read())

    print(f"Uploading note for subject: {subject_key}")

    add_teacher_document(path, subject_key, file.filename)

    notes_db.setdefault(subject_key, []).append({
        "title": file.filename,
        "url": f"http://127.0.0.1:8000/uploads/{filename}"
    })

    with open(NOTES_FILE, "w") as f:
        json.dump(notes_db, f)

    return {"success": True, "message": "Uploaded + learned by AI"}


@app.delete("/notes/{subject}/{filename}")
async def delete_note(subject: str, filename: str):
    subject = subject.lower().strip()

    subject_map = {
        "operating systems": "os",
        "os": "os",
        "computer organization": "coa",
        "coa": "coa",
        "dbms": "dbms",
        "database": "dbms",
        "compiler design": "cd",
        "cd": "cd",
        "python": "python",
        "py": "python"
    }

    subject_key = subject_map.get(subject, subject)

    # Remove from uploads folder
    file_path = None
    for note in notes_db.get(subject_key, []):
        if filename in note["url"]:
            file_path = note["url"].split("/uploads/")[-1]
            break

    if file_path:
        full_path = os.path.join(UPLOAD_DIR, file_path)
        if os.path.exists(full_path):
            os.remove(full_path)

    # Remove from notes_db memory
    notes_db[subject_key] = [
        n for n in notes_db.get(subject_key, [])
        if filename not in n["url"]
    ]

    with open(NOTES_FILE, "w") as f:
        json.dump(notes_db, f)


    return {"success": True, "message": "Note deleted"}


@app.get("/notes/{subject}")
def get_notes(subject: str):

    subject = subject.lower().strip()

    subject_map = {
        "operating systems": "os",
        "os": "os",
        "computer organization": "coa",
        "coa": "coa",
        "dbms": "dbms",
        "database": "dbms",
        "compiler design": "cd",
        "cd": "cd",
        "python": "python",
        "py": "python"
    }

    subject_key = subject_map.get(subject, subject)

    return {"notes": notes_db.get(subject_key, [])}

# ================================
# ROOT
# ================================
@app.get("/")
def root():
    return {"message":"IntelliBot FINAL running"}

# ================================
# Startup load syllabus
# ================================
for subject, path in JSON_FILES.items():
    if os.path.exists(path):
        try:
            build_vectorstore_for_subject(subject, path)
        except Exception as e:
            print(f"❌ Failed loading {subject}: {e}")

# ================================
# RUN SERVER
# ================================
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("server:app", host="127.0.0.1", port=8000, reload=True)
